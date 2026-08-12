"""
智汇中枢 - 标准化工具接口
"""

import os
import json
from datetime import datetime


def search_knowledge_base(db, query, scope=None, scope_id=None, top_k=5, enable_multi_hop=False):
    """跨知识库全文语义搜索"""
    results = []
    with db.session() as s:
        from app.database import File, Folder
        files = s.query(File).join(Folder).filter(
            File.is_deleted == False, File.status == "parsed", Folder.is_deleted == False
        )
        if scope_id:
            if scope == "folder": files = files.filter(File.folder_id == scope_id)
            elif scope == "knowledge_base": files = files.filter(Folder.knowledge_base_id == scope_id)
        files = files.order_by(File.quality_score.desc()).limit(50).all()
    for f in files:
        content = f.content_text or ""
        if query.lower() in content.lower() or query.lower() in f.original_name.lower():
            idx = content.lower().find(query.lower())
            snip = f"...{content[max(0,idx-100):min(len(content),idx+len(query)+200)]}..." if idx >= 0 else content[:300]
            results.append({"file_id": f.id, "file_name": f.original_name, "snippet": snip, "score": f.quality_score or 50, "source": f.original_name})
            if len(results) >= top_k: break
    return {"status": "success", "results": results, "total": len(results), "query": query}


def list_files(db, knowledge_base_id=None, folder_id=None, recursive=False):
    """列出文件"""
    try:
        with db.session() as s:
            from app.database import File, Folder
            q = s.query(File).filter(File.is_deleted == False)
            if folder_id: q = q.filter(File.folder_id == folder_id)
            elif knowledge_base_id: q = q.join(Folder).filter(Folder.knowledge_base_id == knowledge_base_id)
            files = q.order_by(File.updated_at.desc()).limit(100).all()
        return {"status": "success", "files": [{"id": f.id, "name": f.original_name, "type": f.file_type, "size": f.file_size, "status": f.status} for f in files], "total": len(files)}
    except Exception as e:
        return {"status": "error", "files": [], "total": 0, "error": str(e)}


def get_file_content(db, file_id, mode="full", page_start=None, page_end=None):
    """获取文件内容"""
    try:
        with db.session() as s:
            from app.database import File
            f = s.query(File).filter(File.id == file_id).first()
            if not f: return {"status": "error", "content": "文件不存在"}
            content = f.content_text or ""
            if mode == "tables": content = "\n".join(content.split("\n")[:50])
            elif mode == "keypoints": content = content[:500]
            return {"status": "success", "file_id": file_id, "file_name": f.original_name, "content": content[:10000], "total_length": len(f.content_text or "")}
    except Exception as e:
        return {"status": "error", "content": str(e)}


def analyze_data(data_source, operation, filters=None):
    return {"status": "success", "operation": operation, "summary": {"rows": 0, "columns": []}, "details": []}


def predict(data_source, target, period, dimension=None, confidence=0.95):
    return {"status": "success", "target": target, "predictions": [], "confidence": confidence}


def generate_chart(data, chart_type, title, x_label=None, y_label=None, color_scheme="default"):
    """生成图表（支持多种类型 + 地图）"""
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        plt.rcParams["font.sans-serif"] = ["Arial Unicode MS", "SimHei", "PingFang SC", "DejaVu Sans"]
        plt.rcParams["axes.unicode_minus"] = False

        values = data.get("values", [])
        labels = data.get("labels", [])
        cs = ["#1a73e8","#34a853","#ea4335","#fbbc04","#ff6d01","#46bdc6","#7b1fa2","#e91e63","#00bcd4","#8bc34a"]

        fig, ax = plt.subplots(figsize=(12, 7))

        # ── 柱状图（含横向/堆叠） ──
        if chart_type == "bar":
            bars = ax.bar(range(len(values)), values, color=cs[:len(values)], width=0.6, edgecolor="white", linewidth=0.5)
            for bar, v in zip(bars, values):
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01, str(v), ha="center", va="bottom", fontsize=11, fontweight="bold")
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=10)

        elif chart_type == "hbar":
            bars = ax.barh(range(len(values)), values, color=cs[:len(values)], height=0.5, edgecolor="white", linewidth=0.5)
            for bar, v in zip(bars, values):
                ax.text(bar.get_width()+max(values)*0.005, bar.get_y()+bar.get_height()/2, str(v), ha="left", va="center", fontsize=11, fontweight="bold")
            ax.set_yticks(range(len(labels))); ax.set_yticklabels(labels, fontsize=10)
            ax.invert_yaxis()

        elif chart_type == "stacked_bar":
            categories = data.get("categories", labels)
            groups = data.get("groups", ["A","B"])
            group_data = data.get("group_values", [values, [v*0.6 for v in values]])
            bottom = np.zeros(len(categories))
            for i, g in enumerate(group_data[:len(cs)]):
                bars = ax.bar(range(len(categories)), g[:len(categories)], 0.6, bottom=bottom, label=groups[i] if i < len(groups) else f"组{i+1}", color=cs[i % len(cs)])
                bottom += np.array(g[:len(categories)])
            ax.set_xticks(range(len(categories))); ax.set_xticklabels(categories, rotation=25, ha="right", fontsize=10)
            ax.legend(fontsize=11)

        # ── 折线图 ──
        elif chart_type == "line":
            ax.plot(range(len(values)), values, marker="o", linewidth=2.5, color="#1a73e8", markersize=8)
            for i, v in enumerate(values):
                ax.text(i, v+max(values)*0.02, str(v), ha="center", fontsize=11, fontweight="bold")
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=10)

        # ── 饼图/环形图 ──
        elif chart_type == "pie":
            wedges, texts, autotexts = ax.pie(values, labels=labels, autopct="%1.1f%%", colors=cs[:len(values)], startangle=90)
            for at in autotexts: at.set_fontsize(11)

        elif chart_type == "donut":
            wedges, texts, autotexts = ax.pie(values, labels=labels, autopct="%1.1f%%", colors=cs[:len(values)], startangle=90, wedgeprops=dict(width=0.4))
            for at in autotexts: at.set_fontsize(11)

        # ── 散点图/气泡图 ──
        elif chart_type == "scatter":
            ax.scatter(range(len(values)), values, s=120, color="#1a73e8", zorder=5)
            for i, v in enumerate(values): ax.text(i, v+max(values)*0.02, str(v), ha="center", fontsize=10)

        elif chart_type == "bubble":
            sizes = data.get("sizes", [v*10 for v in values])
            ax.scatter(range(len(values)), values, s=[max(s,5) for s in sizes], color=cs[:len(values)], alpha=0.6, zorder=5)
            for i, v in enumerate(values): ax.text(i, v+max(values)*0.02, str(v), ha="center", fontsize=10)
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=10)

        # ── 雷达图 ──
        elif chart_type == "radar":
            angles = np.linspace(0, 2*np.pi, len(values), endpoint=False).tolist()
            values += values[:1]; angles += angles[:1]
            ax = fig.add_subplot(111, polar=True)
            ax.plot(angles, values, "o-", linewidth=2, color="#1a73e8")
            ax.fill(angles, values, alpha=0.1, color="#1a73e8")
            ax.set_xticks(angles[:-1]); ax.set_xticklabels(labels, fontsize=10)

        # ── 箱线图 ──
        elif chart_type == "box":
            box_data = data.get("box_data", [values])
            bp = ax.boxplot(box_data, patch_artist=True, widths=0.5)
            for patch, color in zip(bp["boxes"], cs[:len(box_data)]):
                patch.set_facecolor(color); patch.set_alpha(0.6)
            ax.set_xticklabels(labels if labels else [f"组{i+1}" for i in range(len(box_data))], fontsize=10)

        # ── 热力图 ──
        elif chart_type == "heatmap":
            matrix = data.get("matrix", np.array(values).reshape(-1, len(labels)) if values else np.random.rand(5,5))
            im = ax.imshow(matrix, cmap="Blues", aspect="auto")
            plt.colorbar(im, ax=ax, shrink=0.8)
            rows = data.get("row_labels", [f"行{i+1}" for i in range(len(matrix))])
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, fontsize=9, rotation=25, ha="right")
            ax.set_yticks(range(len(rows))); ax.set_yticklabels(rows, fontsize=9)
            for i in range(len(matrix)):
                for j in range(len(matrix[0])):
                    ax.text(j, i, f"{matrix[i][j]:.1f}", ha="center", va="center", fontsize=9, color="white" if matrix[i][j] > matrix.max()/2 else "black")

        # ── 地图（Cartopy 真实地图） ──
        elif chart_type == "map":
            try:
                import cartopy.crs as ccrs
                import cartopy.feature as cfeature
                fig = plt.figure(figsize=(14, 10))
                ax = fig.add_subplot(1, 1, 1, projection=ccrs.PlateCarree())
                ax.set_extent([73, 135, 18, 54], crs=ccrs.PlateCarree())
                ax.add_feature(cfeature.LAND, facecolor="#f0f0f0")
                ax.add_feature(cfeature.OCEAN, facecolor="#dce8f5")
                ax.add_feature(cfeature.COASTLINE, linewidth=0.8)
                ax.add_feature(cfeature.BORDERS, linewidth=0.5, linestyle=":", edgecolor="#888")
                ax.add_feature(cfeature.LAKES, alpha=0.5)
                ax.add_feature(cfeature.RIVERS, alpha=0.3)
                ax.gridlines(draw_labels=True, linestyle="--", alpha=0.3)
                lons = data.get("longitudes", [116.4, 121.5, 114.1, 113.3, 120.2])
                lats = data.get("latitudes", [39.9, 31.2, 22.5, 23.1, 30.2])
                cities = data.get("cities", labels) or ["北京","上海","广州","深圳","杭州"]
                vals = values if values else [100]*len(cities)
                sizes = [max(v, 5) for v in vals]
                sc = ax.scatter(lons, lats, s=[s*4 for s in sizes], c="#ea4335", alpha=0.7, transform=ccrs.PlateCarree(), zorder=5, edgecolors="white", linewidth=0.5)
                for i, city in enumerate(cities):
                    ax.text(lons[i]+0.5, lats[i]+0.5, f"{city}\n{vals[i]}", fontsize=10, transform=ccrs.PlateCarree(), fontweight="bold")
                plt.title(title, fontsize=18, pad=20, fontweight="bold")
            except ImportError:
                plt.close()
                return {"status": "error", "error": "cartopy未安装，请运行: pip install cartopy"}
        else:
            ax.set_title(title, fontsize=18, pad=24, fontweight="bold")
            if x_label: ax.set_xlabel(x_label, fontsize=13)
            if y_label: ax.set_ylabel(y_label, fontsize=13)
            if chart_type not in ("pie", "donut", "radar", "heatmap"):
                ax.grid(axis="y", alpha=0.3, linestyle="--"); ax.set_axisbelow(True)

        import tempfile
        img_path = os.path.join(tempfile.gettempdir(), f"chart_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
        plt.savefig(img_path, dpi=200, bbox_inches="tight", facecolor="white", edgecolor="none")
        plt.close()
        return {"status": "success", "chart_type": chart_type, "title": title, "image_path": img_path}
    except Exception as e:
        import traceback
        return {"status": "error", "error": str(e), "trace": traceback.format_exc()[:300]}


def generate_report(title, sources, template=None, sections=None, tone="professional", include_appendix=False):
    """生成报告"""
    content_parts = [f"# {title}\n"]
    if sections:
        for s in sections:
            if isinstance(s, dict):
                section_title = s.get("title", s.get("heading", "章节"))
                section_content = s.get("content", s.get("text", ""))
                content_parts.append(f"\n## {section_title}\n")
                if section_content:
                    content_parts.append(f"{section_content}\n")
                else:
                    content_parts.append("（分析内容）\n")
            else:
                content_parts.append(f"\n## {s}\n\n（分析内容）\n")
    return {"status": "success", "title": title, "format": "markdown", "content": "\n".join(content_parts), "file_path": None}


def generate_ppt(title, sources, theme=None, slide_count=None, include_notes=False):
    """生成 PPT（支持多主题设计）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn

        # 主题配色
        THEMES = {
            "professional": {  # 蓝白商务
                "primary": RGBColor(0x1a, 0x73, 0xe8),
                "secondary": RGBColor(0x34, 0xa8, 0x53),
                "accent": RGBColor(0xea, 0x43, 0x35),
                "bg": RGBColor(0xf8, 0xf9, 0xfa),
                "text": RGBColor(0x1a, 0x1a, 0x1a),
                "light_text": RGBColor(0x66, 0x66, 0x66),
                "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
            },
            "modern": {  # 深色炫酷
                "primary": RGBColor(0x66, 0x77, 0xFF),
                "secondary": RGBColor(0x00, 0xC9, 0xA7),
                "accent": RGBColor(0xFF, 0x6B, 0x6B),
                "bg": RGBColor(0x1a, 0x1a, 0x2e),
                "text": RGBColor(0xFF, 0xFF, 0xFF),
                "light_text": RGBColor(0xAA, 0xAA, 0xCC),
                "card_bg": RGBColor(0x25, 0x25, 0x40),
            },
            "elegant": {  # 莫兰迪风格
                "primary": RGBColor(0x7C, 0x6F, 0x9A),
                "secondary": RGBColor(0xB8, 0xA9, 0xC9),
                "accent": RGBColor(0xD4, 0xA5, 0x7A),
                "bg": RGBColor(0xFA, 0xF5, 0xF0),
                "text": RGBColor(0x3D, 0x34, 0x4D),
                "light_text": RGBColor(0x8E, 0x83, 0x9C),
                "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
            },
            "vibrant": {  # 鲜艳活力
                "primary": RGBColor(0xFF, 0x6B, 0x35),
                "secondary": RGBColor(0xFF, 0xA8, 0x00),
                "accent": RGBColor(0x00, 0xB4, 0xD8),
                "bg": RGBColor(0xFF, 0xF8, 0xF0),
                "text": RGBColor(0x2D, 0x2D, 0x2D),
                "light_text": RGBColor(0x88, 0x88, 0x88),
                "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
            },
            "dark": {  # 极简暗色
                "primary": RGBColor(0x00, 0xD2, 0xFF),
                "secondary": RGBColor(0x6C, 0x63, 0xFF),
                "accent": RGBColor(0xFF, 0xD6, 0x00),
                "bg": RGBColor(0x0D, 0x0D, 0x0D),
                "text": RGBColor(0xEE, 0xEE, 0xEE),
                "light_text": RGBColor(0x88, 0x88, 0x88),
                "card_bg": RGBColor(0x1A, 0x1A, 0x1A),
            },
        }
        t = THEMES.get(theme or "professional", THEMES["professional"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        # ── 封面 ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = t["primary"]
        # 装饰条纹
        for i in range(3):
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(5.5 + i * 0.6), Inches(13.333), Pt(2)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = t["secondary"]
            stripe.line.fill.background()
            stripe.fill.fore_color.brightness = 0.0
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # ── 内容页 ──
        for item in sources:
            if isinstance(item, dict):
                slide_title = item.get("title", "")
                bullets = item.get("content") or item.get("bullets") or []

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = t["bg"]

                # 顶部色条
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), Inches(13.333), Inches(0.08)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = t["primary"]
                bar.line.fill.background()

                # 标题
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = t["primary"]

                # 标题下分隔线
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(1.35), Inches(2), Pt(3)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = t["secondary"]
                line.line.fill.background()

                # 内容（每个要点一个卡片）
                if bullets:
                    card_y = Inches(1.7)
                    card_width = Inches(11.5)
                    for i, b in enumerate(bullets):
                        card = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), card_y + Inches(i * 0.85),
                            card_width, Inches(0.7)
                        )
                        card.fill.solid()
                        card.fill.fore_color.rgb = t["card_bg"]
                        card.line.color.rgb = t["primary"]
                        card.line.width = Pt(0.5)
                        # 设置阴影效果（通过 XML）
                        card.shadow.inherit = False
                        # 内容文字
                        txBox2 = slide.shapes.add_textbox(
                            Inches(1.2), card_y + Inches(i * 0.85) + Inches(0.05),
                            card_width - Inches(0.8), Inches(0.6)
                        )
                        tf2 = txBox2.text_frame
                        tf2.word_wrap = True
                        p2 = tf2.paragraphs[0]
                        p2.text = str(b)
                        p2.font.size = Pt(16)
                        p2.font.color.rgb = t["text"]
                        p2.space_after = Pt(4)

                # 页码
                page_num = sources.index(item) + 2 if item in sources else 2
                num_box = slide.shapes.add_textbox(
                    Inches(12), Inches(7.0), Inches(1), Inches(0.4)
                )
                np = num_box.text_frame.paragraphs[0]
                np.text = str(page_num)
                np.font.size = Pt(11)
                np.font.color.rgb = t["light_text"]
                np.alignment = PP_ALIGN.RIGHT
            else:
                # 纯文字页
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(item)[:200]
                p.font.size = Pt(20)
                p.font.color.rgb = t["text"]

        import tempfile
        filename = f"ppt_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = os.path.join(tempfile.gettempdir(), filename)
        prs.save(file_path)
        return {"status": "success", "title": title, "slide_count": len(sources) + 1, "file_path": file_path, "theme": theme or "professional"}
    except Exception as e:
        import traceback
        return {"status": "error", "title": title, "slide_count": 0, "file_path": None, "error": str(e)}


def export_result(result_ref, format="md", filename=None, merge=False, export_dir=None):
    """导出结果"""
    if not filename: filename = f"export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{format}"
    file_path = os.path.join(export_dir or "/tmp", filename) if filename else None
    return {"status": "success", "format": format, "filename": filename, "file_path": file_path}
