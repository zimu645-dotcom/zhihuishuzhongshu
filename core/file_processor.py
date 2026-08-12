"""
智汇中枢 - 文件解析处理器
支持 TXT/MD/PDF/DOCX/XLSX/PPTX/CSV 等格式自动解析
"""

import os
import re
import hashlib
from datetime import datetime
from typing import Optional


def parse_file(file_path: str) -> dict:
    """
    解析文件，返回结构化内容

    Returns:
        {
            "success": bool,
            "content": str,          # 解析出的纯文本
            "page_count": int,       # 页数（文档类）
            "tables": list,          # 表格数据（Excel类）
            "quality_score": float,  # 质量评分 0-100
            "error": str             # 错误信息
        }
    """
    ext = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".txt": _parse_text,
        ".md": _parse_text,
        ".json": _parse_text,
        ".csv": _parse_csv,
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pptx": _parse_pptx,
        ".png": _parse_image,
        ".jpg": _parse_image,
        ".jpeg": _parse_image,
        ".gif": _parse_image,
        ".bmp": _parse_image,
        ".webp": _parse_image,
    }

    parser = parsers.get(ext)
    if not parser:
        return {
            "success": False,
            "content": f"不支持的文件格式: {ext}",
            "page_count": 0,
            "tables": [],
            "quality_score": 0,
            "error": f"unsupported_format: {ext}"
        }

    try:
        result = parser(file_path)
        # 计算质量评分
        result["quality_score"] = _calculate_quality_score(result["content"])
        return result
    except Exception as e:
        return {
            "success": False,
            "content": f"解析失败: {str(e)}",
            "page_count": 0,
            "tables": [],
            "quality_score": 0,
            "error": str(e)
        }


def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> list:
    """将文本切分成块"""
    if not text:
        return []

    # 先按段落分割
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current_chunk = ""
    current_len = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        para_len = len(para)
        if current_len + para_len <= chunk_size:
            current_chunk += para + "\n\n"
            current_len += para_len
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            # 如果段落本身超长，按句子分割
            if para_len > chunk_size:
                sentences = re.split(r'(?<=[。！？.!?])\s*', para)
                temp_chunk = ""
                temp_len = 0
                for sent in sentences:
                    sent_len = len(sent)
                    if temp_len + sent_len <= chunk_size:
                        temp_chunk += sent
                        temp_len += sent_len
                    else:
                        if temp_chunk:
                            chunks.append(temp_chunk.strip())
                        temp_chunk = sent
                        temp_len = sent_len
                if temp_chunk:
                    chunks.append(temp_chunk.strip())
            else:
                current_chunk = para + "\n\n"
                current_len = para_len

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks


def _calculate_quality_score(content: str) -> float:
    """计算文本质量评分 0-100"""
    if not content:
        return 0

    score = 100
    text_len = len(content)

    # 长度扣分
    if text_len < 10:
        score -= 80
    elif text_len < 50:
        score -= 40
    elif text_len < 200:
        score -= 10

    # 乱码检测（中英文混合率异常）
    ascii_chars = sum(1 for c in content if ord(c) < 128)
    total_chars = len(content)
    if total_chars > 0:
        ascii_ratio = ascii_chars / total_chars
        if ascii_ratio > 0.95 and ascii_ratio < 1.0:
            score -= 20  # 可能是二进制文件被当文本读
        if ascii_ratio < 0.1:
            score -= 30  # 可能是纯二进制

    # 重复内容检测
    lines = content.split('\n')
    if len(lines) > 5:
        unique_ratio = len(set(lines)) / len(lines)
        if unique_ratio < 0.3:
            score -= 30  # 大量重复行

    return max(0, min(100, score))


def _parse_text(file_path: str) -> dict:
    """解析纯文本文件"""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return {
        "success": True,
        "content": content,
        "page_count": max(1, len(content) // 2000 + 1),
        "tables": [],
        "error": None
    }


def _parse_csv(file_path: str) -> dict:
    """解析 CSV 文件"""
    import csv
    import io

    content_parts = []
    tables = []
    total_rows = 0

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            content_parts.append(" | ".join(row))
            total_rows += 1
            if total_rows <= 50:
                tables.append(row)

    content = "\n".join(content_parts)
    return {
        "success": True,
        "content": content,
        "page_count": max(1, total_rows // 50 + 1),
        "tables": tables,
        "error": None
    }


def _parse_pdf(file_path: str) -> dict:
    """解析 PDF 文件"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return _fallback_parse(file_path, "PDF")

    doc = fitz.open(file_path)
    content_parts = []
    page_count = len(doc)

    for page_num in range(page_count):
        page = doc[page_num]
        text = page.get_text()
        if text.strip():
            content_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")

    doc.close()
    return {
        "success": True,
        "content": "\n\n".join(content_parts),
        "page_count": page_count,
        "tables": [],
        "error": None
    }


def _parse_docx(file_path: str) -> dict:
    """解析 Word 文档"""
    try:
        from docx import Document
    except ImportError:
        return _fallback_parse(file_path, "Word")

    doc = Document(file_path)
    content_parts = []
    tables = []

    for para in doc.paragraphs:
        if para.text.strip():
            content_parts.append(para.text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            content_parts.append(" | ".join(cells))
            rows.append(cells)

    return {
        "success": True,
        "content": "\n\n".join(content_parts),
        "page_count": max(1, len(content_parts) // 50 + 1),
        "tables": tables[:10],
        "error": None
    }


def _parse_xlsx(file_path: str) -> dict:
    """解析 Excel 文件"""
    try:
        import openpyxl
    except ImportError:
        return _fallback_parse(file_path, "Excel")

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    content_parts = []
    tables = []
    total_rows = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        content_parts.append(f"\n=== 工作表: {sheet_name} ===")

        sheet_rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                content_parts.append(row_text)
                sheet_rows.append(row)
                total_rows += 1
                if total_rows <= 100:
                    tables.append(row)

    wb.close()
    return {
        "success": True,
        "content": "\n".join(content_parts),
        "page_count": max(1, total_rows // 50 + 1),
        "tables": tables,
        "error": None
    }


def _parse_pptx(file_path: str) -> dict:
    """解析 PPT 文件"""
    try:
        from pptx import Presentation
    except ImportError:
        return _fallback_parse(file_path, "PPT")

    prs = Presentation(file_path)
    content_parts = []
    tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        content_parts.append(f"\n--- 幻灯片 {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content_parts.append(shape.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    content_parts.append(" | ".join(cells))

    if not content_parts:
        content_parts.append("(PPT 无文本内容，可能全是图片)")

    return {
        "success": True,
        "content": "\n".join(content_parts),
        "page_count": len(prs.slides),
        "tables": [],
        "error": None
    }


def _parse_image(file_path: str) -> dict:
    """解析图片 - 提取基本信息"""
    try:
        from PIL import Image
    except ImportError:
        return _fallback_parse(file_path, "图片")

    img = Image.open(file_path)
    width, height = img.size
    mode = img.mode
    fmt = img.format

    # 尝试 OCR (如果安装了 pytesseract)
    ocr_text = ""
    try:
        import pytesseract
        ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
    except Exception:
        pass

    content = (
        f"图片信息:\n"
        f"格式: {fmt}\n"
        f"尺寸: {width} x {height} 像素\n"
        f"色彩模式: {mode}\n"
        f"文件: {os.path.basename(file_path)}\n"
    )
    if ocr_text.strip():
        content += f"\nOCR 识别文字:\n{ocr_text.strip()}"

    return {
        "success": True,
        "content": content,
        "page_count": 1,
        "tables": [],
        "error": None
    }


def _fallback_parse(file_path: str, file_type: str) -> dict:
    """兜底解析 - 尝试按文本读取"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read(10000)
        return {
            "success": True if content.strip() else False,
            "content": content if content.strip() else f"({file_type} 文件，内容为空或无法读取)",
            "page_count": 1,
            "tables": [],
            "error": None if content.strip() else "empty_content"
        }
    except Exception as e:
        return {
            "success": False,
            "content": f"({file_type} 文件，无法解析文本内容)",
            "page_count": 1,
            "tables": [],
            "quality_score": 0,
            "error": str(e)
        }
